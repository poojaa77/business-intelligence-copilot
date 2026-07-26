"""
report_export.py
Builds downloadable PowerPoint (.pptx) and PDF reports summarizing the
KPI dashboard, trends, forecasts, anomalies, and the AI diagnostic.
Chart images are rendered from the same Plotly figures shown on screen
(via kaleido) so the export matches what the user sees.
"""

import io
from datetime import date


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def fig_to_png_bytes(fig, width=900, height=500):
    return fig.to_image(format="png", width=width, height=height, scale=2)


# ---------------------------------------------------------------------------
# PowerPoint export
# ---------------------------------------------------------------------------

def build_pptx(title: str, diagnostic: dict, chart_figs: dict, forecasts: dict, anomalies: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # --- Title slide ---
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Business Intelligence Copilot — generated {date.today().isoformat()}"

    # --- Executive summary slide ---
    slide = prs.slides.add_slide(blank_layout)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
    box.text_frame.text = "Executive Summary"
    box.text_frame.paragraphs[0].font.size = Pt(28)
    box.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.7))
    tf = body.text_frame
    tf.word_wrap = True
    tf.text = diagnostic.get("executive_summary", "")
    tf.paragraphs[0].font.size = Pt(18)

    def _bullet_slide(header, items):
        s = prs.slides.add_slide(blank_layout)
        b = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        b.text_frame.text = header
        b.text_frame.paragraphs[0].font.size = Pt(26)
        b.text_frame.paragraphs[0].font.bold = True

        content = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.7))
        ctf = content.text_frame
        ctf.word_wrap = True
        if not items:
            ctf.text = "None identified."
        else:
            ctf.text = f"• {items[0]}"
            ctf.paragraphs[0].font.size = Pt(16)
            for item in items[1:]:
                p = ctf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(16)
        return s

    _bullet_slide("Positive Trends", diagnostic.get("positive_trends", []))
    _bullet_slide("Areas of Concern", diagnostic.get("areas_of_concern", []))
    _bullet_slide("Risks", diagnostic.get("risks", []))
    _bullet_slide("Growth Opportunities", diagnostic.get("growth_opportunities", []))
    _bullet_slide("Recommendations", diagnostic.get("recommendations", []))

    # --- Forecast & anomaly summary slide ---
    fc_lines = [
        f"{m}: forecast {f['forecast']} (confidence: {f['confidence']})"
        for m, f in forecasts.items()
    ]
    _bullet_slide("Next-Month Forecasts", fc_lines)

    an_lines = []
    for metric, found in anomalies.items():
        for a in found:
            an_lines.append(f"{metric} — {a['period']}: {a['direction']} to {a['value']}")
    _bullet_slide("Detected Anomalies", an_lines)

    # --- Chart slides ---
    for metric_name, fig in chart_figs.items():
        s = prs.slides.add_slide(blank_layout)
        header = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
        header.text_frame.text = metric_name
        header.text_frame.paragraphs[0].font.size = Pt(24)
        header.text_frame.paragraphs[0].font.bold = True
        try:
            img_bytes = fig_to_png_bytes(fig)
            s.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.8), Inches(1.2), width=Inches(11.5))
        except Exception:
            note = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(1))
            note.text_frame.text = "(Chart image unavailable — install the 'kaleido' package to enable chart export.)"

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ---------------------------------------------------------------------------
# PDF export
# ---------------------------------------------------------------------------

def build_pdf(title: str, diagnostic: dict, chart_figs: dict, forecasts: dict, anomalies: dict) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image, ListFlowable, ListItem
    )

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=letter,
                             topMargin=0.6 * inch, bottomMargin=0.6 * inch)
    styles = getSampleStyleSheet()
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    body = styles["BodyText"]

    story = [
        Paragraph(title, h1),
        Paragraph(f"Generated {date.today().isoformat()} — Business Intelligence Copilot", body),
        Spacer(1, 16),
        Paragraph("Executive Summary", h2),
        Paragraph(diagnostic.get("executive_summary", ""), body),
        Spacer(1, 12),
    ]

    def _section(header, items):
        story.append(Paragraph(header, h2))
        if items:
            story.append(ListFlowable(
                [ListItem(Paragraph(i, body)) for i in items], bulletType="bullet"
            ))
        else:
            story.append(Paragraph("None identified.", body))
        story.append(Spacer(1, 12))

    _section("Positive Trends", diagnostic.get("positive_trends", []))
    _section("Areas of Concern", diagnostic.get("areas_of_concern", []))
    _section("Risks", diagnostic.get("risks", []))
    _section("Growth Opportunities", diagnostic.get("growth_opportunities", []))
    _section("Recommendations", diagnostic.get("recommendations", []))

    _section("Next-Month Forecasts", [
        f"{m}: forecast {f['forecast']} (confidence: {f['confidence']})"
        for m, f in forecasts.items()
    ])

    an_lines = []
    for metric, found in anomalies.items():
        for a in found:
            an_lines.append(f"{metric} — {a['period']}: {a['direction']} to {a['value']}")
    _section("Detected Anomalies", an_lines)

    story.append(Paragraph("Charts", h2))
    for metric_name, fig in chart_figs.items():
        try:
            img_bytes = fig_to_png_bytes(fig, width=650, height=380)
            story.append(Paragraph(metric_name, styles["Heading3"]))
            story.append(Image(io.BytesIO(img_bytes), width=6.2 * inch, height=3.6 * inch))
            story.append(Spacer(1, 10))
        except Exception:
            story.append(Paragraph(
                f"{metric_name}: (chart image unavailable — install 'kaleido')", body
            ))

    doc.build(story)
    return output.getvalue()
